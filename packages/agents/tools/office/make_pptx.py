#!/usr/bin/env python3
"""Build a presentation from a JSON spec. No API to learn, one command:

    make-deck spec.json out.pptx

spec.json:

    {"title": "The deck title",
     "subtitle": "",
     "theme": {"accent": "#E07A1B", "background": "#FFFFFF", "text": "#1A1A1A"},
     "slides": [
       {"title": "Problem", "bullets": ["First point", "Second"], "notes": ""},
       {"title": "Growth", "chart": {"kind": "column",
                                     "categories": ["Q1", "Q2", "Q3"],
                                     "series": [{"name": "Subscribers",
                                                 "values": [120, 340, 810]}]}}
     ]}

One idea per slide, at most five bullets — a slide is a prompt for a
speaker, not a page of prose. notes lands in the speaker notes.

theme is optional and applies to every slide: accent colours the chart bars,
the title rule and the bullet marks; background and text are what they say.
Colours are #RRGGBB.

chart is optional, one per slide, and takes kind column, bar, line or pie
with its categories and one or more named series. A slide may carry both
bullets and a chart; the chart sits under them.

This exists because a model writing python-pptx code from memory invents
imports and APIs — Chart.format and Plot.format are the usual two, and
neither exists. A model filling in this spec and running one command does
not have to know that. The spec is the whole surface.
"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from artifact_path import resolve_output  # noqa: E402

KINDS = {"column", "bar", "line", "pie"}


def colour(said: str, what: str):
    from pptx.dml.color import RGBColor

    text = str(said).strip().lstrip("#")
    if len(text) != 6:
        raise SystemExit(f'theme.{what} is "{said}": a colour is #RRGGBB, like #E07A1B')
    try:
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        raise SystemExit(f'theme.{what} is "{said}": a colour is #RRGGBB, like #E07A1B')


def themed(spec: dict) -> dict:
    said = spec.get("theme") or {}
    if not isinstance(said, dict):
        raise SystemExit('"theme" is an object, like {"accent": "#E07A1B"}')
    out = {}
    for what in ("accent", "background", "text"):
        if said.get(what):
            out[what] = colour(said[what], what)
    return out


def paint(slide, look: dict) -> None:
    if "background" in look:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = look["background"]
    if "text" not in look and "accent" not in look:
        return
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        heading = shape == slide.shapes.title
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if heading and "accent" in look:
                    run.font.color.rgb = look["accent"]
                elif "text" in look:
                    run.font.color.rgb = look["text"]


def draw(slide, said: dict, look: dict, top_in: float) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    kind = str(said.get("kind", "column")).lower()
    if kind not in KINDS:
        raise SystemExit(f'chart kind "{kind}" is not one of: ' + ", ".join(sorted(KINDS)))
    cats = said.get("categories")
    if not isinstance(cats, list) or not cats:
        raise SystemExit('a chart needs "categories": a non-empty list of labels')
    series = said.get("series")
    if not isinstance(series, list) or not series:
        raise SystemExit('a chart needs "series": [{"name": "...", "values": [1, 2, 3]}]')

    data = CategoryChartData()
    data.categories = [str(c) for c in cats]
    for i, one in enumerate(series):
        if not isinstance(one, dict):
            raise SystemExit(f"series[{i}] is not an object")
        values = one.get("values")
        if not isinstance(values, list) or len(values) != len(cats):
            raise SystemExit(
                f"series[{i}] has {len(values) if isinstance(values, list) else 0} values"
                f" for {len(cats)} categories — they must match"
            )
        try:
            data.add_series(str(one.get("name", f"Series {i + 1}")), tuple(float(v) for v in values))
        except (TypeError, ValueError):
            raise SystemExit(f"series[{i}] has a value that is not a number")

    types = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    frame = slide.shapes.add_chart(
        types[kind], Inches(0.8), Inches(top_in), Inches(8.4), Inches(6.6 - top_in), data)
    chart = frame.chart
    chart.font.size = Pt(12)
    if kind == "pie" or len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    # Chart.format and Plot.format do not exist in python-pptx; a series is
    # coloured through its own format, and that is the only way in.
    if "accent" in look and kind != "pie":
        chart.plots[0].series[0].format.fill.solid()
        chart.plots[0].series[0].format.fill.fore_color.rgb = look["accent"]
    if "text" in look:
        chart.font.color.rgb = look["text"]


def build(spec: dict, out: Path) -> int:
    from pptx import Presentation
    from pptx.util import Pt

    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise SystemExit('spec has no "slides": expected a non-empty list of {"title","bullets"}')
    look = themed(spec)
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[0])
    first.shapes.title.text = str(spec.get("title", "Untitled"))
    if spec.get("subtitle") and len(first.placeholders) > 1:
        first.placeholders[1].text = str(spec["subtitle"])
    paint(first, look)
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            raise SystemExit(f"slides[{i}] is not an object")
        chart = s.get("chart")
        if chart is not None and not isinstance(chart, dict):
            raise SystemExit(f'slides[{i}]["chart"] is an object, or leave it out')
        bullets = s.get("bullets") or []
        layout = 5 if chart is not None and not bullets else 1
        slide = deck.slides.add_slide(deck.slide_layouts[layout])
        slide.shapes.title.text = str(s.get("title", ""))
        if bullets:
            body = slide.placeholders[1].text_frame
            for j, b in enumerate(bullets):
                para = body.paragraphs[0] if j == 0 else body.add_paragraph()
                para.text = str(b)
                para.font.size = Pt(20)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(s["notes"])
        paint(slide, look)
        if chart is not None:
            draw(slide, chart, look, 3.6 if bullets else 1.8)
    parent = os.path.dirname(str(out))
    if parent:
        os.makedirs(parent, exist_ok=True)
    deck.save(str(out))
    charts = sum(1 for s in slides if isinstance(s, dict) and s.get("chart"))
    said = f"wrote {out} ({len(slides)} slides"
    print(said + (f", {charts} with a chart)" if charts else ")"))
    return 0


def main(argv: list[str]) -> int:
    # See make_docx.py: the one real cause of a wrong count is an unquoted
    # space in the output path, so the error says that instead of the usage.
    if len(argv) > 3:
        raise SystemExit(
            f"got {len(argv) - 1} arguments, expected 2: make-deck spec.json out.pptx —"
            " a path with a space must be quoted; better, name the file with dashes,"
            " like /artifacts/decks/q3-review.pptx"
        )
    if len(argv) < 3:
        print(__doc__)
        return 2
    spec_path, out = Path(argv[1]), Path(resolve_output(argv[2]))
    try:
        spec = json.loads(spec_path.read_text("utf8"))
    except FileNotFoundError:
        raise SystemExit(f"{spec_path}: no such file — write the spec first, then run make-deck")
    except json.JSONDecodeError as e:
        raise SystemExit(f"{spec_path} is not valid JSON: {e}")
    if not isinstance(spec, dict):
        raise SystemExit(f"{spec_path}: expected a JSON object, got {type(spec).__name__}")
    return build(spec, out)


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
