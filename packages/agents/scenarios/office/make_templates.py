# Build the real office templates the capability pages start from.
#
# A template is a document, not a prompt: styled headings, tables that already
# have their columns, slide layouts already chosen. The agent opens it and
# fills it in, which is a different job from writing one from nothing — and
# the reason a person picks a template at all.
#
# Placeholders are <ANGLED> so they are findable by the editor skill and
# obviously unfinished to a reader who opens the file before the agent does.
import os
from docx import Document
from docx.shared import Pt, RGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Pt as PPt

OUT = os.path.dirname(os.path.abspath(__file__))

def status_report():
    d = Document()
    d.add_heading("<TITLE>", level=0)
    p = d.add_paragraph("<PERIOD> · <AUTHOR>")
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    for heading, body in [
        ("Context", "<one paragraph: what the reader needs before the detail>"),
        ("What happened", None),
        ("What is next", None),
        ("Decisions waiting", None),
    ]:
        d.add_heading(heading, level=1)
        if body:
            d.add_paragraph(body)
        else:
            for _ in range(3):
                d.add_paragraph("<item>", style="List Bullet")
    d.add_heading("Numbers", level=1)
    t = d.add_table(rows=3, cols=3)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Measure", "This period", "Last period"]):
        cell = t.cell(0, i)
        cell.text = h
        cell.paragraphs[0].runs[0].font.bold = True
    for r in range(1, 3):
        for c in range(3):
            t.cell(r, c).text = "<value>"
    d.save(os.path.join(OUT, "status-report.docx"))

def budget_tracker():
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    head = ["Item", "Planned", "Actual", "Variance"]
    ws.append(head)
    fill = PatternFill("solid", fgColor="EEEEEE")
    for c in ws[1]:
        c.font = Font(bold=True)
        c.fill = fill
    for r in range(2, 12):
        ws.cell(r, 1, "<item>")
        ws.cell(r, 2, 0)
        ws.cell(r, 3, 0)
        # The variance column is a formula, which is the point of shipping a
        # real workbook: the agent fills items and numbers, the sheet does the
        # arithmetic and keeps doing it after the reader edits a cell.
        ws.cell(r, 4, f"=C{r}-B{r}")
    ws.cell(12, 1, "Total").font = Font(bold=True)
    for col in "BCD":
        ws[f"{col}12"] = f"=SUM({col}2:{col}11)"
        ws[f"{col}12"].font = Font(bold=True)
    for i, w in enumerate([34, 14, 14, 14], start=1):
        ws.column_dimensions[get_column_letter(i)].width = w
    wb.save(os.path.join(OUT, "budget-tracker.xlsx"))

def pitch_deck():
    deck = Presentation()
    title = deck.slides.add_slide(deck.slide_layouts[0])
    title.shapes.title.text = "<COMPANY>"
    title.placeholders[1].text = "<the one sentence>"
    for name, hint in [
        ("Problem", "<whose, how often, what it costs them today>"),
        ("Insight", "<what you know that the room does not>"),
        ("What we built", "<the thing, not the roadmap>"),
        ("Why it wins", "<the moat, the wedge, or the unfair advantage>"),
        ("The ask", "<amount, hire, decision — be specific>"),
    ]:
        s = deck.slides.add_slide(deck.slide_layouts[1])
        s.shapes.title.text = name
        tf = s.placeholders[1].text_frame
        tf.paragraphs[0].text = hint
        tf.paragraphs[0].font.size = PPt(20)
        for _ in range(2):
            para = tf.add_paragraph()
            para.text = "<point>"
            para.font.size = PPt(20)
        s.notes_slide.notes_text_frame.text = "<what the speaker says here>"
    deck.save(os.path.join(OUT, "pitch-deck.pptx"))

# --- four more of each, so a capability page is a shelf rather than one card ---

def meeting_notes():
    d = Document()
    d.add_heading("<MEETING>", level=0)
    p = d.add_paragraph("<DATE> · <WHO CALLED IT>")
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    d.add_heading("Attendees", level=1)
    for _ in range(3):
        d.add_paragraph("<name — role>", style="List Bullet")
    d.add_heading("Decisions", level=1)
    for _ in range(2):
        d.add_paragraph("<what was decided, and why now>", style="List Bullet")
    d.add_heading("Actions", level=1)
    t = d.add_table(rows=4, cols=3)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Action", "Owner", "Due"]):
        c = t.cell(0, i); c.text = h; c.paragraphs[0].runs[0].font.bold = True
    for r in range(1, 4):
        for c, v in enumerate(["<action>", "<who>", "<when>"]):
            t.cell(r, c).text = v
    d.add_heading("Parked", level=1)
    d.add_paragraph("<raised, not resolved — where it goes next>", style="List Bullet")
    d.save(os.path.join(OUT, "meeting-notes.docx"))

def project_brief():
    d = Document()
    d.add_heading("<PROJECT>", level=0)
    p = d.add_paragraph("<OWNER> · <TARGET DATE>")
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    for heading, hint in [
        ("Problem", "<who hurts, how often, what it costs them>"),
        ("Goal", "<the one sentence this ships against>"),
        ("Not doing", None),
        ("Risks", None),
    ]:
        d.add_heading(heading, level=1)
        if hint:
            d.add_paragraph(hint)
        else:
            for _ in range(3):
                d.add_paragraph("<item>", style="List Bullet")
    d.add_heading("Milestones", level=1)
    t = d.add_table(rows=4, cols=3)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Milestone", "Date", "Proof it landed"]):
        c = t.cell(0, i); c.text = h; c.paragraphs[0].runs[0].font.bold = True
    for r in range(1, 4):
        for c, v in enumerate(["<milestone>", "<date>", "<what a reader can check>"]):
            t.cell(r, c).text = v
    d.save(os.path.join(OUT, "project-brief.docx"))

def proposal():
    d = Document()
    d.add_heading("<PROPOSAL>", level=0)
    p = d.add_paragraph("For <AUDIENCE> · <DATE>")
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    d.add_heading("Summary", level=1)
    d.add_paragraph("<the recommendation in three sentences, written last>")
    d.add_heading("Options", level=1)
    t = d.add_table(rows=4, cols=4)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["Option", "Cost", "Risk", "Why / why not"]):
        c = t.cell(0, i); c.text = h; c.paragraphs[0].runs[0].font.bold = True
    for r in range(1, 4):
        for c, v in enumerate(["<option>", "<cost>", "<risk>", "<the honest line>"]):
            t.cell(r, c).text = v
    d.add_heading("Recommendation", level=1)
    d.add_paragraph("<which option, and the two reasons that decided it>")
    d.add_heading("What happens next", level=1)
    for _ in range(3):
        d.add_paragraph("<step — owner — date>", style="List Bullet")
    d.save(os.path.join(OUT, "proposal.docx"))

def decision_memo():
    d = Document()
    d.add_heading("<DECISION>", level=0)
    p = d.add_paragraph("<STATUS: proposed | decided> · <OWNER> · <DATE>")
    p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    for heading, hint in [
        ("Context", "<what forced a decision at all>"),
        ("The decision", "<one sentence, active voice>"),
        ("Options considered", None),
        ("Consequences", "<what gets harder, who is affected, what we stop doing>"),
        ("Revisit when", "<the observable fact that reopens this>"),
    ]:
        d.add_heading(heading, level=1)
        if hint:
            d.add_paragraph(hint)
        else:
            for _ in range(3):
                d.add_paragraph("<option — why it lost>", style="List Bullet")
    d.save(os.path.join(OUT, "decision-memo.docx"))

def project_timeline():
    wb = Workbook()
    ws = wb.active
    ws.title = "Timeline"
    ws.append(["Task", "Owner", "Start", "End", "Days", "Status"])
    fill = PatternFill("solid", fgColor="EEEEEE")
    for c in ws[1]:
        c.font = Font(bold=True); c.fill = fill
    for r in range(2, 14):
        ws.cell(r, 1, "<task>"); ws.cell(r, 2, "<who>")
        ws.cell(r, 3, "<start>"); ws.cell(r, 4, "<end>")
        ws.cell(r, 5, f"=D{r}-C{r}")
        ws.cell(r, 6, "<planned | doing | done>")
    for i, w in enumerate([36, 14, 12, 12, 8, 18], start=1):
        ws.column_dimensions[get_column_letter(i)].width = w
    wb.save(os.path.join(OUT, "project-timeline.xlsx"))

def invoice():
    wb = Workbook()
    ws = wb.active
    ws.title = "Invoice"
    ws["A1"] = "<YOUR COMPANY>"; ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "<address · email>"
    ws["A4"] = "Bill to:"; ws["A4"].font = Font(bold=True)
    ws["A5"] = "<client>"; ws["A6"] = "<their address>"
    ws["D4"] = "Invoice #"; ws["E4"] = "<number>"
    ws["D5"] = "Date"; ws["E5"] = "<date>"
    ws["D6"] = "Due"; ws["E6"] = "<date + terms>"
    head_row = 8
    for i, h in enumerate(["Description", "Qty", "Unit price", "Amount"], start=1):
        c = ws.cell(head_row, i, h); c.font = Font(bold=True)
        c.fill = PatternFill("solid", fgColor="EEEEEE")
    for r in range(head_row + 1, head_row + 7):
        ws.cell(r, 1, "<line item>"); ws.cell(r, 2, 0); ws.cell(r, 3, 0)
        ws.cell(r, 4, f"=B{r}*C{r}")
    total = head_row + 7
    ws.cell(total, 3, "Total").font = Font(bold=True)
    ws.cell(total, 4, f"=SUM(D{head_row+1}:D{total-1})").font = Font(bold=True)
    for i, w in enumerate([42, 8, 12, 12], start=1):
        ws.column_dimensions[get_column_letter(i)].width = w
    wb.save(os.path.join(OUT, "invoice.xlsx"))

def expense_report():
    wb = Workbook()
    ws = wb.active
    ws.title = "Expenses"
    ws.append(["Date", "Category", "Description", "Amount", "Receipt?"])
    fill = PatternFill("solid", fgColor="EEEEEE")
    for c in ws[1]:
        c.font = Font(bold=True); c.fill = fill
    for r in range(2, 14):
        ws.cell(r, 1, "<date>"); ws.cell(r, 2, "<travel | meals | tools | other>")
        ws.cell(r, 3, "<what and why>"); ws.cell(r, 4, 0); ws.cell(r, 5, "<y/n>")
    ws.cell(14, 3, "Total").font = Font(bold=True)
    ws.cell(14, 4, "=SUM(D2:D13)").font = Font(bold=True)
    for i, w in enumerate([12, 20, 40, 12, 10], start=1):
        ws.column_dimensions[get_column_letter(i)].width = w
    wb.save(os.path.join(OUT, "expense-report.xlsx"))

def content_calendar():
    wb = Workbook()
    ws = wb.active
    ws.title = "Calendar"
    ws.append(["Publish date", "Channel", "Title", "Owner", "Status", "Link"])
    fill = PatternFill("solid", fgColor="EEEEEE")
    for c in ws[1]:
        c.font = Font(bold=True); c.fill = fill
    for r in range(2, 14):
        ws.cell(r, 1, "<date>"); ws.cell(r, 2, "<blog | social | email>")
        ws.cell(r, 3, "<working title>"); ws.cell(r, 4, "<who>")
        ws.cell(r, 5, "<idea | drafting | review | live>"); ws.cell(r, 6, "<url>")
    for i, w in enumerate([14, 16, 40, 14, 16, 28], start=1):
        ws.column_dimensions[get_column_letter(i)].width = w
    wb.save(os.path.join(OUT, "content-calendar.xlsx"))

def _deck(path, title, subtitle, slides):
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[0])
    first.shapes.title.text = title
    first.placeholders[1].text = subtitle
    for name, hint in slides:
        s = deck.slides.add_slide(deck.slide_layouts[1])
        s.shapes.title.text = name
        tf = s.placeholders[1].text_frame
        tf.paragraphs[0].text = hint
        tf.paragraphs[0].font.size = PPt(20)
        for _ in range(2):
            para = tf.add_paragraph(); para.text = "<point>"; para.font.size = PPt(20)
        s.notes_slide.notes_text_frame.text = "<what the speaker says here>"
    deck.save(os.path.join(OUT, path))

def kickoff_deck():
    _deck("project-kickoff.pptx", "<PROJECT>", "<one line: what ships, for whom>", [
        ("Why now", "<the fact that made this worth starting this quarter>"),
        ("Goal", "<what done looks like, in numbers a reader can check>"),
        ("Scope", "<in — and the two things loudly out>"),
        ("Plan", "<the three milestones, with dates>"),
        ("Team", "<who owns what, one name per line>"),
        ("Risks", "<the one that keeps you up, and the mitigation>"),
    ])

def quarterly_review():
    _deck("quarterly-review.pptx", "<TEAM> — <QUARTER>", "<the quarter in one sentence>", [
        ("Highlights", "<the three things worth the room's time>"),
        ("Numbers", "<targets against actuals — no adjectives>"),
        ("What missed", "<what, why, and what changed because of it>"),
        ("Learned", "<the thing you would tell last quarter's self>"),
        ("Next quarter", "<the bets, ranked, with owners>"),
    ])

def team_update():
    _deck("team-update.pptx", "<TEAM> update", "<date — cadence>", [
        ("Wins", "<shipped, closed, unblocked — since last time>"),
        ("Metrics", "<the two or three the team steers by>"),
        ("In flight", "<what lands next, and when>"),
        ("Help wanted", "<the specific ask, addressed to someone in the room>"),
    ])

def roadmap_deck():
    _deck("product-roadmap.pptx", "<PRODUCT> roadmap", "<horizon — e.g. the next two quarters>", [
        ("Themes", "<the two or three problems this roadmap is about>"),
        ("Now", "<building today — lands within the quarter>"),
        ("Next", "<committed, not started — the following quarter>"),
        ("Later", "<believed in, not committed — no dates on purpose>"),
        ("Tradeoffs", "<what this roadmap says no to>"),
    ])

status_report(); budget_tracker(); pitch_deck()
meeting_notes(); project_brief(); proposal(); decision_memo()
project_timeline(); invoice(); expense_report(); content_calendar()
kickoff_deck(); quarterly_review(); team_update(); roadmap_deck()
print("built:", sorted(f for f in os.listdir(OUT) if f.split('.')[-1] in ('docx','xlsx','pptx')))
