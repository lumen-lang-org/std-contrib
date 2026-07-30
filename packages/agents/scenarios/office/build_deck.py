# The make-deck builder.
#
# spec = {
#   "path": "/artifacts/decks/pitch.pptx",
#   "title": "Lumen",                      # title slide
#   "subtitle": "",                        # optional
#   "slides": [
#     { "title": "Problem", "bullets": ["...", "..."], "notes": "" },
#   ],
# }
import os
from pptx import Presentation
from pptx.util import Pt

def build(spec):
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[0])
    first.shapes.title.text = spec.get("title", "Untitled")
    if spec.get("subtitle") and len(first.placeholders) > 1:
        first.placeholders[1].text = spec["subtitle"]
    for s in spec.get("slides", []):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1].text_frame
        for i, b in enumerate(s.get("bullets", [])):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = b
            para.font.size = Pt(20)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    path = spec["path"]
    os.makedirs(os.path.dirname(path), exist_ok=True)
    deck.save(path)
    print("wrote", path)
