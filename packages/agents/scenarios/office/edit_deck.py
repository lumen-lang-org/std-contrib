# Fill a deck that already exists.
#
# The template chose its layouts; this replaces the text inside them. A deck
# generated beside the template would use python-pptx's defaults instead —
# which is a different deck, whatever the words say.
import os
from pptx import Presentation


def fill(path, fills, slides=None):
    """Replace <PLACEHOLDER> text everywhere, then optionally rewrite bullets.

    slides = { "Problem": ["...", "..."] } replaces the body of the slide
    whose title matches, keeping that slide's layout and its notes.
    """
    deck = Presentation(path)
    hits = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    out = run.text
                    for key, value in fills.items():
                        out = out.replace(key, str(value))
                    if out != run.text:
                        run.text = out
                        hits += 1

    for title, bullets in (slides or {}).items():
        for slide in deck.slides:
            if not slide.shapes.title or slide.shapes.title.text.strip() != title:
                continue
            body = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    body = shape.text_frame
                    break
            if body is None:
                continue
            # The first paragraph carries the layout's own formatting, so it
            # is rewritten rather than cleared and rebuilt.
            size = body.paragraphs[0].font.size
            body.paragraphs[0].text = bullets[0] if bullets else ""
            body.paragraphs[0].font.size = size
            for extra in body.paragraphs[1:]:
                extra._p.getparent().remove(extra._p)
            for text in bullets[1:]:
                para = body.add_paragraph()
                para.text = text
                para.font.size = size

    deck.save(path)
    print(f"filled {hits} placeholders in {os.path.basename(path)}")
    return hits
